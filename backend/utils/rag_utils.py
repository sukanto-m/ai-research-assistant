import os
from langchain.document_loaders import PyPDFLoader
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.embeddings import OpenAIEmbeddings
from langchain.vectorstores import FAISS
from langchain.schema import Document
from docx import Document as DocxDocument
from pptx import Presentation
import openpyxl
import pandas as pd
from bs4 import BeautifulSoup

embedding_model = OpenAIEmbeddings(model="text-embedding-3-small")
vector_store = None

def extract_text_from_docx(file_path: str) -> str:
    doc = DocxDocument(file_path)
    return '\n'.join([para.text for para in doc.paragraphs])

def extract_text_from_pptx(file_path: str) -> str:
    prs = Presentation(file_path)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return '\n'.join(text)

def extract_text_from_xlsx(file_path: str) -> str:
    wb = openpyxl.load_workbook(file_path)
    text = []
    for sheet in wb.worksheets:
        for row in sheet.iter_rows():
            for cell in row:
                if cell.value:
                    text.append(str(cell.value))
    return '\n'.join(text)

def extract_text_from_txt(file_path: str) -> str:
    with open(file_path, "r", encoding="utf-8") as f:
        return f.read()

def extract_text_from_md(file_path: str) -> str:
    with open(file_path, "r", encoding="utf-8") as f:
        return f.read()

def extract_text_from_csv(file_path: str) -> str:
    df = pd.read_csv(file_path)
    return df.to_string(index=False)

def extract_text_from_html(file_path: str) -> str:
    with open(file_path, "r", encoding="utf-8") as f:
        soup = BeautifulSoup(f.read(), "html.parser")
    return soup.get_text()

def load_and_embed_file(file_path: str):
    global vector_store

    ext = os.path.splitext(file_path)[1].lower()

    if ext == ".pdf":
        loader = PyPDFLoader(file_path)
        documents = loader.load()
    elif ext == ".docx":
        content = extract_text_from_docx(file_path)
        documents = [Document(page_content=content)]
    elif ext == ".pptx":
        content = extract_text_from_pptx(file_path)
        documents = [Document(page_content=content)]
    elif ext == ".xlsx":
        content = extract_text_from_xlsx(file_path)
        documents = [Document(page_content=content)]
    elif ext == ".txt":
        content = extract_text_from_txt(file_path)
        documents = [Document(page_content=content)]
    elif ext == ".md":
        content = extract_text_from_md(file_path)
        documents = [Document(page_content=content)]
    elif ext == ".csv":
        content = extract_text_from_csv(file_path)
        documents = [Document(page_content=content)]
    elif ext == ".html":
        content = extract_text_from_html(file_path)
        documents = [Document(page_content=content)]
    else:
        raise ValueError("Unsupported file format")

    splitter = RecursiveCharacterTextSplitter(chunk_size=500, chunk_overlap=50)
    chunks = splitter.split_documents(documents)
    vector_store = FAISS.from_documents(chunks, embedding_model)
    print(f"✅ Document indexed with {len(chunks)} chunks.")

def get_relevant_chunks(query: str, k: int = 3) -> str:
    if vector_store is None:
        return "⚠️ No document uploaded yet."

    docs = vector_store.similarity_search(query, k=k)
    return "\n\n".join([doc.page_content for doc in docs])
